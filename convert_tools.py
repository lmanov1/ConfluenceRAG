from docx import Document
import mimetypes
import pytesseract

# Example conversion
from pdf2image import convert_from_path
import pdf2image

from PIL import Image
from pdfplumber import open as pdf_open
from pptx import Presentation

import logging


def check_pdf2image():
    pdf_path = "Mbuf_debug_tool.pdf"
    # Set the path to Poppler's `pdftoppm` executable
    images = convert_from_path(pdf_path)
    for i, image in enumerate(images):
        image.save(f"page_{i+1}.jpg", "JPEG")
        #image.show()
        image.close()
    
def check_tesseract():
    # Create a simple image with text
    image = Image.new("RGB", (300, 100), color=(255, 255, 255))
    #image.show()
    image.save("test_image.png")
    image.close()
    pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"

    # Extract text from the image
    try:
        text = pytesseract.image_to_string("test_image.png")
        print("Tesseract is working. Extracted text:", text)
    except pytesseract.TesseractNotFoundError:
        print("Tesseract is not installed or not found.")
    
    
# Function to extract text from .docx files
def extract_text_from_docx(file_path):
    text = ""
    document = Document(file_path)
    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pdf(file_path):
    text = ""
    with pdf_open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_pptx(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_image(file_path):
    image = Image.open(file_path)
    
    pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"
    text = pytesseract.image_to_string(image)
    return text


# Function to process non-text documents
def process_non_text_document(file_path):
    mime_type, _ = mimetypes.guess_type(file_path)
    print(f"Going to scrape {file_path} with mime type: {mime_type} ")
    if mime_type == "application/pdf":
        return extract_text_from_pdf(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or mime_type == "application/msword":
        return extract_text_from_docx(file_path)
    elif mime_type == "application/vnd.ms-powerpoint" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file_path)
    elif mime_type and mime_type.startswith("image/"):
        return extract_text_from_image(file_path)
    else:
        return None

# Process documents from Confluence
def preprocess_documents(documents, chunk_size=500):
    logging.basicConfig(level=logging.DEBUG)
    chunks = []
    for doc in documents:
        if hasattr(doc, 'file_path') and doc.file_path:  # Non-text document
            text = process_non_text_document(doc.file_path)
        else:  # Text content
            text = doc.page_content
        if text:
            chunks.extend([text[i:i+chunk_size] for i in range(0, len(text), chunk_size)])
    return chunks

# check_tesseract()
# check_pdf2image()
